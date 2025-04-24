from PyPDF2 import PdfReader
from pptx import Presentation

def extract_pdf_text(filepath):
    reader = PdfReader(filepath)
    text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
    return text



def extract_ppt_text(filepath: str) -> str:
    prs = Presentation(filepath)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)